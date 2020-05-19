# -*- coding: utf-8 -*-
"""
Created on Sat May  9 22:52:43 2020

@author: Ramy.Saad
"""

pip install python-pptx

import pandas as pd
from datetime import date
from pptx import Presentation
import matplotlib.pyplot as plt
import seaborn as sns
import matplotlib.image as mpimg 
from PIL import Image 

#Input filelocation followed by sheetname 

dffile = "C:/Users/ramy.saad/Desktop/Solutions Python Scripts/Mizkan/Data/Ragu Performance March 2020.xlsx"
df = pd.read_excel(dffile, 'Sheet1') 
df = df[['Placement','Creative Version','Metric Date','Week','Total Impressions','Unique Clicks','Measured','Viewed']]

def genAudience(s):
    if 'category_v' in s:
        return 'Category'
    if '_category' in s:
        return 'Category'
    if 'competitive_v' in s:
        return 'Competitive'
    if 'd_competitive' in s:
        return 'Competitive'
    if 'previous_v' in s:
        return 'Past'
    if '_previous' in s:
        return 'Past'
    if 'Category_O' in s:
        return 'Category_OWS'
    if 'Contextual_O' in s:
        return 'Contextual_OWS'
    if 'Previous_O' in s:
        return 'Previous_OWS'
    elif 'Competitive_O' in s:
        return 'Competitive_OWS'
    return ''

df['Audience'] = df['Creative Version'].apply(genAudience)
df = df[df.Audience != ''] #delete rows without 

#make this into function! + loops + add in begining! 
df3=pd.pivot_table(df, ['Measured','Viewed','Unique Clicks','Total Impressions'], ['Audience'])
df3['CTR']= (df3['Unique Clicks']/df3['Total Impressions'])*100
df3["IVR"]= (df3["Viewed"]/df3["Measured"])*100
df3 = df3[['CTR','IVR']]
df3 = df3.reset_index()
df3.sort_values('IVR', inplace=True)
df3.columns = ['Audience','CTR (%)','IVR (%)']

#Create combo chart
fig = plt.figure()
fig, ax1 = plt.subplots(figsize=(12,6))
color = 'tab:green'
#bar plot creation
ax1.set_title('IVR and CTR by Audience', fontsize=16)
ax1.set_xlabel('Audience', fontsize=16)
ax1.set_ylabel('CTR (%)', fontsize=16)
ax1 = sns.barplot(x='Audience', y='CTR (%)', data = df3, palette='winter')
ax1.tick_params(axis='y')
#specify we want to share the same x-axis
ax2 = ax1.twinx()
color = 'tab:red'
#line plot creation
ax2.set_ylabel('IVR by Audience', fontsize=16)
ax2 = sns.lineplot(x='Audience', y='IVR (%)', data = df3, sort=False, color=color)
ax2.tick_params(axis='y', color=color)
#show plot
plt.show()
fig.savefig('C:/Users/ramy.saad/Desktop/Solutions Python Scripts/Mizkan/Output/new-file-name.png')

#prs = Presentation('C:/Users/ramy.saad/Desktop/Solutions Python Scripts/Mizkan/Ragu DCO Results FY2020 March_2020.pptx')
#prs.save('C:/Users/ramy.saad/Desktop/Solutions Python Scripts/Mizkan/Output/new-file-name.pptx')

def add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]
 
    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size
 
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width
 
    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)
 
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
 
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side

topIVR = df3['IVR (%)'].max()
i = df3['IVR (%)'].idxmax()
topIVRAudience = df3['Audience'][i]

topCTR = df3['CTR (%)'].max()
i = df3['CTR (%)'].idxmax()
topCTRAudience = df3['Audience'][i]
topCTR = round(topCTR,2)
topIVR = round(topIVR,2)
results = "Highest CTR audience category ="+ ' ' + topCTRAudience + ' ' + str(topCTR) +'%' + '\nHighest CTR audience category ='  + ' ' + topIVRAudience + ' ' + str(topIVR)+'%'
        
prs = Presentation('C:/Users/ramy.saad/Desktop/Solutions Python Scripts/Mizkan/Ragu DCO Results FY2020 March_2020 2.pptx')
#prs = Presentation()
layout = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout)
title = slide.shapes.title
title.text = "CTV and IVR by Audience"
subtitle = slide.placeholders[2]
subtitle.text = results
add_image(slide,1,"C:/Users/ramy.saad/Desktop/Solutions Python Scripts/Mizkan/Output/new-file-name.png")
prs.save('C:/Users/ramy.saad/Desktop/Solutions Python Scripts/Mizkan/Output/Mizkan.pptx')


